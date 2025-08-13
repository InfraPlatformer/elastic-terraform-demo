#!/usr/bin/env python3
"""
PowerPoint Presentation Generator
Converts markdown presentation to PowerPoint slides
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re

def create_title_slide(prs, title, subtitle=""):
    """Create a title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    if subtitle:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
    
    return slide

def create_content_slide(prs, title, content):
    """Create a content slide with title and bullet points"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    content_shape = slide.placeholders[1]
    content_shape.text = content
    
    return slide

def create_code_slide(prs, title, code, language="bash"):
    """Create a slide with code block"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Add code as text box
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = code
    
    # Format as code (monospace font)
    p = tf.paragraphs[0]
    p.font.name = 'Courier New'
    p.font.size = Pt(10)
    
    return slide

def create_table_slide(prs, title, headers, rows):
    """Create a slide with a table"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Create table
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)
    
    table = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table
    
    # Add headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
    
    # Add rows
    for i, row in enumerate(rows):
        for j, cell_text in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(cell_text)
    
    return slide

def parse_markdown_presentation(md_file):
    """Parse markdown presentation file"""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    slides = []
    slide_sections = content.split('---')
    
    for section in slide_sections:
        if not section.strip():
            continue
            
        lines = section.strip().split('\n')
        if not lines:
            continue
            
        # Extract slide title
        title = ""
        content_lines = []
        code_block = ""
        in_code_block = False
        table_data = None
        
        for line in lines:
            if line.startswith('## Slide'):
                # Extract title from slide header
                title_match = re.search(r'## Slide \d+: (.+)', line)
                if title_match:
                    title = title_match.group(1)
            elif line.startswith('```'):
                in_code_block = not in_code_block
                if in_code_block:
                    code_block = ""
                else:
                    content_lines.append(f"```\n{code_block}\n```")
            elif in_code_block:
                code_block += line + '\n'
            elif line.startswith('|') and '|' in line:
                # Table data
                if table_data is None:
                    table_data = {'headers': [], 'rows': []}
                cells = [cell.strip() for cell in line.split('|')[1:-1]]
                if not table_data['headers']:
                    table_data['headers'] = cells
                else:
                    table_data['rows'].append(cells)
            elif line.strip() and not line.startswith('#'):
                content_lines.append(line)
        
        if title:
            slides.append({
                'title': title,
                'content': '\n'.join(content_lines),
                'table_data': table_data
            })
    
    return slides

def create_presentation_from_markdown(md_file, output_file):
    """Create PowerPoint presentation from markdown file"""
    # Create presentation
    prs = Presentation()
    
    # Parse markdown
    slides_data = parse_markdown_presentation(md_file)
    
    for slide_data in slides_data:
        title = slide_data['title']
        content = slide_data['content']
        table_data = slide_data['table_data']
        
        if table_data:
            create_table_slide(prs, title, table_data['headers'], table_data['rows'])
        elif '```' in content:
            # Extract code block
            code_match = re.search(r'```(?:\w+)?\n(.*?)\n```', content, re.DOTALL)
            if code_match:
                code = code_match.group(1)
                create_code_slide(prs, title, code)
            else:
                create_content_slide(prs, title, content)
        else:
            create_content_slide(prs, title, content)
    
    # Save presentation
    prs.save(output_file)
    print(f"Presentation saved as: {output_file}")

def main():
    """Main function"""
    # Check if python-pptx is installed
    try:
        import pptx
    except ImportError:
        print("Error: python-pptx is not installed.")
        print("Install it with: pip install python-pptx")
        sys.exit(1)
    
    # File paths
    md_file = "../docs/Meetup_Presentation.md"
    output_file = "../docs/Elastic_Terraform_Presentation.pptx"
    
    # Check if markdown file exists
    if not os.path.exists(md_file):
        print(f"Error: Markdown file not found: {md_file}")
        sys.exit(1)
    
    # Create presentation
    try:
        create_presentation_from_markdown(md_file, output_file)
        print("✅ Presentation created successfully!")
        print(f"📁 Location: {output_file}")
    except Exception as e:
        print(f"Error creating presentation: {e}")
        sys.exit(1)

if __name__ == "__main__":
    main()
